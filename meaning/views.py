from django.shortcuts import render

# Create your views here.
from django.shortcuts import render
from django.contrib.auth.models import User
from django.contrib import messages
from django.contrib.auth import authenticate,login,logout
from django.http import HttpResponse
from math import ceil
import docx2txt
from PyDictionary import PyDictionary
import PyPDF2
from PyPDF2 import PdfFileReader, PdfFileWriter
import pdfplumber
from pptx import Presentation
import googletrans
from googletrans import Translator
from .models import Contact
import re
import os
import gtts
from playsound import playsound
# Create your views here.
def index(request):
    return render(request,"meaning/index.html")

def about(request):
    return render(request,"meaning/about.html")

def Contact(request):
    if request.method=="POST":
        print(request)
        name=request.POST.get('name','')
        email=request.POST.get('email','')
        phone=request.POST.get('phone','')
        desc=request.POST.get('desc','')
        contact=Contact(name=name,email=email,phone=phone,desc=desc)
        contact.save()
    return render(request,"meaning/contact.html")

def notes(request):
    return render(request,"meaning/notes.html")


def pro(request):
    if request.method=="POST":
        text=request.POST.get('text5','')
        lang=request.POST.get('text6','')
        tts = gtts.gTTS(text, lang=lang)
        tts.save("hello.mp3")
        playsound("hello.mp3")
        os.remove("hello.mp3")
    return render(request,"meaning/index.html")

def trans(request):
    if request.method=="POST":
        text=request.POST.get('text','')
        lang=request.POST.get('lang','')
        translator = Translator(service_urls=['translate.googleapis.com'])
        result1=translator.translate(text, dest=lang)
        final={}
        final={lang:result1.text}

            
    return render(request,"meaning/index.html",{'final':final})

def result(request):
    if request.method=="POST":
        set1={'1','2','3','4','5','6','7','8','9','0'}
        print(request)
        file1=request.POST.get('file1','')
        file2=request.POST.get('file2','')
        file3=request.POST.get('file3','')
        file4=request.POST.get('file4','')
        file5=request.POST.get('file5','')
        file6=request.POST.get('file6','')
        file7=request.POST.get('file7','')
        if("docx" in file1):
            my_text = docx2txt.process(f"D:\projectfiles\{file1}")
            my_text = re.sub(r'[^\w\s]', '', my_text) 
            men_text=my_text.split()
            dictionary=PyDictionary()
            final={}
            for ele in men_text:
                if(len(ele)>5):
                    if(ele in final):
                        continue
                    elif(len(list(set1.intersection(set(ele))))>0):
                        continue
                    final[ele]=dictionary.meaning(ele)

            print(final)
        elif("pdf" in file2):
            pdf = PyPDF2.PdfFileReader(open(f"D:\projectfiles\{file2}", "rb"))
            my_text1=""
            for page in pdf.pages:
                my_text1+=page.extractText()
            my_text1 = re.sub(r'[^\w\s]', '', my_text1)
            men_text1=my_text1.split()
            dictionary=PyDictionary()
            final={}
            for ele in men_text1:
                if(len(ele)>5):
                    if(ele in final):
                        continue
                    elif(len(list(set1.intersection(set(ele))))>0):
                        continue
                    else:
                        final[ele]=dictionary.meaning(ele)
        elif("ppt" in file3):
            print("hiii")
            prs = Presentation(f"D:/projectfiles/{file3}")

            string=""
            for slide in prs.slides:
                for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            string=string+shape.text
                            
                            
            res = re.sub(r'[^\w\s]', '', string)

            men_text=res.split()
        
            dictionary=PyDictionary()
            final={}
            for ele in men_text:

                if(len(ele)>5):
                    if(ele in final):
                        continue
                    elif(len(list(set1.intersection(set(ele))))>0):
                        continue
                    else:
                        final[ele]=dictionary.meaning(ele)
        elif("docx" in file4):
            my_text = docx2txt.process(f"D:\projectfiles\{file4}")
            my_text = re.sub(r'[^\w\s]', '', my_text) 
            men_text=my_text.split()
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            for ele in men_text:
                if(len(ele)>5):
                    if(ele in final):
                        continue
                    elif(len(list(set1.intersection(set(ele))))>0):
                        continue
                    result=translator.translate(ele, dest='hi')
                    final[ele] =result.text

            print(final)
        elif("pdf" in file4):
            my_text=""
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            i=0
            with pdfplumber.open(f"D:\projectfiles\{file4}") as pdf:
                for page in pdf.pages:
                    my_text+=page.extract_text()
            my_text = re.sub(r'[^\w\s]', '', my_text)
            men_text=my_text.split()
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            for ele in men_text:
                if(len(ele)>5):
                    if(ele in final):
                        continue
                    elif(len(list(set1.intersection(set(ele))))>0):
                        continue
                    else:
                        result=translator.translate(ele, dest='hi')
                        final[ele] =result.text
        elif("ppt" in file4):
            print("hiii")
            prs = Presentation(f"D:/projectfiles/{file4}")

            string=""
            for slide in prs.slides:
                for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            string=string+shape.text
                            
                            
            res = re.sub(r'[^\w\s]', '', string)

            men_text=res.split()
        
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            for ele in men_text:

                if(len(ele)>5):
                    if(ele in final):
                        continue
                    elif(len(list(set1.intersection(set(ele))))>0):
                        continue
                    else:
                        result=translator.translate(ele, dest='hi')
                        final[ele] =result.text
        elif("docx" in file5):
            my_text = docx2txt.process(f"D:\projectfiles\{file5}")
            my_text = re.sub(r'[^\w\s]', '', my_text) 
            men_text=my_text.split()
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            for ele in men_text:
                if(len(ele)>5):
                    if(ele in final):
                        continue
                    elif(len(list(set1.intersection(set(ele))))>0):
                        continue

                    result=translator.translate(ele, dest='mr')
                    final[ele] =result.text

            print(final)
        elif("pdf" in file5):
            my_text=""
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            i=0
            with pdfplumber.open(f"D:\projectfiles\{file5}") as pdf:
                for page in pdf.pages:
                    my_text+=page.extract_text()
            my_text = re.sub(r'[^\w\s]', '', my_text)
            men_text=my_text.split()
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            for ele in men_text:
                if(len(ele)>5):
                    if(ele in final):
                        continue
                    elif(len(list(set1.intersection(set(ele))))>0):
                        continue
                    else:
                        result=translator.translate(ele, dest='mr')
                        final[ele] =result.text
        elif("ppt" in file5):
            print("hiii")
            prs = Presentation(f"D:/projectfiles/{file5}")

            string=""
            for slide in prs.slides:
                for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            string=string+shape.text
                            
                            
            res = re.sub(r'[^\w\s]', '', string)

            men_text=res.split()
        
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            for ele in men_text:

                if(len(ele)>5):
                    if(ele in final):
                        continue
                    elif(len(list(set1.intersection(set(ele))))>0):
                        continue
                    else:
                        result=translator.translate(ele, dest='mr')
                        final[ele] =result.text
        elif("docx" in file6):
            my_text = docx2txt.process(f"D:\projectfiles\{file6}")
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            result=translator.translate(my_text, dest='hi')
            final["Translation"] =result.text
        
        elif("pdf" in file6):
            
            my_text=""
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            i=0
            with pdfplumber.open(f"D:\projectfiles\{file6}") as pdf:
                for page in pdf.pages:
                    i+=1
                    my_text=page.extract_text()
                    result=translator.translate(my_text, dest='hi')
                    final[i] =result.text
        
        elif("ppt" in file6):
            prs = Presentation(f"D:/projectfiles/{file7}")
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            i=0
            string=""
            for slide in prs.slides:
                for shape in slide.shapes:
                    i+=1
                    if hasattr(shape, "text"):
                        string=shape.text
                        result=translator.translate(my_text, dest='hi')
                        final[i] =result.text

        elif("docx" in file7):
            my_text = docx2txt.process(f"D:\projectfiles\{file7}")
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            result=translator.translate(my_text, dest='mr')
            final["Translation"] =result.text
        
        elif("pdf" in file7):
            my_text=""
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            i=0
            with pdfplumber.open(f"D:\projectfiles\{file7}") as pdf:
                for page in pdf.pages:
                    i+=1
                    my_text=page.extract_text()
                    print(my_text)
                    result=translator.translate(my_text, dest='mr')
                    final[i] =result.text
        
        elif("ppt" in file7):
            prs = Presentation(f"D:/projectfiles/{file7}")
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            i=0
            string=""
            my_text=""
            for slide in prs.slides:
                for shape in slide.shapes:
                    i+=1
                    if hasattr(shape, "text"):
                        string=shape.text
                        result=translator.translate(my_text, dest='mr')
                        final[i] =result.text
        
        
        else:
            final={"Error": {"Jo Uplode karne ko bola vo kar na":"bheag"}}

            

    return render(request,"meaning/result.html",{"final":final})

def prodectView(request,myid):
    return render(request,"meaning/prodview.html")



